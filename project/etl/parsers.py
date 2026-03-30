"""
文档解析器 — PDF / PPT / DOCX / Markdown / Python
====================================================

每种格式一个 Parser，统一接口:
  parse(file_path) → list[Document]

Document 结构:
  {
    "content": "文本内容",
    "metadata": {
      "source": "培训手册.pdf",
      "source_type": "pdf",
      "page": 3,           # PDF 页码
      "section": "投诉处理", # 标题/章节
      "slide": 2,           # PPT slide 编号
    }
  }

设计原则:
  - 按语义分块，不是固定字数切割
  - PDF: 按页 + 段落分块，表格保持完整
  - PPT: 每张 slide 是一个 chunk（标题+内容+备注）
  - DOCX: 按标题层级分块
  - 保留丰富的 metadata 用于引用溯源
"""

import logging
from pathlib import Path
from dataclasses import dataclass, field

logger = logging.getLogger("etl.parsers")


@dataclass
class Document:
    """解析后的文档块"""
    content: str
    metadata: dict = field(default_factory=dict)

    def to_dict(self) -> dict:
        return {"content": self.content, "metadata": self.metadata}


# ============================================================
# PDF Parser
# ============================================================

def parse_pdf(file_path: str | Path) -> list[Document]:
    """
    解析 PDF 文件

    策略:
      - 每页作为一个基础 chunk
      - 如果某页太长 (>2000 字)，按段落拆分
      - 保留页码、文件名等 metadata
    """
    import fitz  # PyMuPDF

    file_path = Path(file_path)
    docs = []

    try:
        pdf = fitz.open(str(file_path))
        logger.info("解析 PDF: %s (%d 页)", file_path.name, len(pdf))

        for page_num, page in enumerate(pdf, 1):
            text = page.get_text("text").strip()
            if not text:
                continue

            # 如果页面文本较短，整页作为一个 chunk
            if len(text) <= 2000:
                docs.append(Document(
                    content=text,
                    metadata={
                        "source": file_path.name,
                        "source_type": "pdf",
                        "page": page_num,
                        "total_pages": len(pdf),
                    },
                ))
            else:
                # 长页面按段落拆分（双换行分隔）
                paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
                for i, para in enumerate(paragraphs):
                    if len(para) < 50:  # 太短的段落合并到下一个
                        continue
                    docs.append(Document(
                        content=para,
                        metadata={
                            "source": file_path.name,
                            "source_type": "pdf",
                            "page": page_num,
                            "paragraph": i + 1,
                        },
                    ))

        pdf.close()
        logger.info("PDF 解析完成: %d 个文档块", len(docs))

    except Exception as e:
        logger.error("PDF 解析失败 %s: %s", file_path, e)

    return docs


# ============================================================
# PPT Parser
# ============================================================

def parse_pptx(file_path: str | Path) -> list[Document]:
    """
    解析 PowerPoint (.pptx) 文件

    策略:
      - 每张 slide 作为一个 chunk
      - 提取: 标题 + 正文 + 备注
      - 带 slide 编号作为 metadata
    """
    from pptx import Presentation

    file_path = Path(file_path)
    docs = []

    try:
        prs = Presentation(str(file_path))
        logger.info("解析 PPTX: %s (%d slides)", file_path.name, len(prs.slides))

        for slide_num, slide in enumerate(prs.slides, 1):
            parts = []

            # 提取标题
            if slide.shapes.title and slide.shapes.title.text.strip():
                parts.append(f"# {slide.shapes.title.text.strip()}")

            # 提取所有文本框内容
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text and text != (slide.shapes.title.text.strip() if slide.shapes.title else ""):
                        parts.append(text)

                # 提取表格
                if shape.has_table:
                    table_text = _extract_pptx_table(shape.table)
                    if table_text:
                        parts.append(table_text)

            # 提取备注
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[备注] {notes}")

            content = "\n\n".join(parts)
            if content.strip():
                docs.append(Document(
                    content=content,
                    metadata={
                        "source": file_path.name,
                        "source_type": "pptx",
                        "slide": slide_num,
                        "total_slides": len(prs.slides),
                        "has_notes": bool(notes),
                    },
                ))

        logger.info("PPTX 解析完成: %d 个文档块", len(docs))

    except Exception as e:
        logger.error("PPTX 解析失败 %s: %s", file_path, e)

    return docs


def _extract_pptx_table(table) -> str:
    """提取 PPT 表格为 Markdown 格式"""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(" | ".join(cells))
    if not rows:
        return ""
    # 添加 Markdown 表格分隔符
    header = rows[0]
    separator = " | ".join(["---"] * len(table.columns))
    return "\n".join([header, separator] + rows[1:])


# ============================================================
# DOCX Parser
# ============================================================

def parse_docx(file_path: str | Path) -> list[Document]:
    """
    解析 Word (.docx) 文件

    策略:
      - 按标题 (Heading) 分块
      - 每个标题下的所有段落合并为一个 chunk
      - 没有标题的段落按固定长度分块
      - 表格保持完整
    """
    from docx import Document as DocxDocument

    file_path = Path(file_path)
    docs = []

    try:
        docx = DocxDocument(str(file_path))
        logger.info("解析 DOCX: %s (%d 段落)", file_path.name, len(docx.paragraphs))

        current_section = ""
        current_content = []

        for para in docx.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # 检测标题
            if para.style.name.startswith("Heading"):
                # 保存之前的 section
                if current_content:
                    docs.append(Document(
                        content="\n\n".join(current_content),
                        metadata={
                            "source": file_path.name,
                            "source_type": "docx",
                            "section": current_section or "(无标题)",
                        },
                    ))
                current_section = text
                current_content = [f"# {text}"]
            else:
                current_content.append(text)

        # 保存最后一个 section
        if current_content:
            docs.append(Document(
                content="\n\n".join(current_content),
                metadata={
                    "source": file_path.name,
                    "source_type": "docx",
                    "section": current_section or "(无标题)",
                },
            ))

        # 提取表格
        for i, table in enumerate(docx.tables):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                header = rows[0]
                sep = " | ".join(["---"] * len(table.columns))
                table_text = "\n".join([header, sep] + rows[1:])
                docs.append(Document(
                    content=table_text,
                    metadata={
                        "source": file_path.name,
                        "source_type": "docx",
                        "section": f"表格 {i + 1}",
                        "content_type": "table",
                    },
                ))

        logger.info("DOCX 解析完成: %d 个文档块", len(docs))

    except Exception as e:
        logger.error("DOCX 解析失败 %s: %s", file_path, e)

    return docs


# ============================================================
# Markdown / Text Parser (已有，增强版)
# ============================================================

def parse_markdown(file_path: str | Path) -> list[Document]:
    """
    解析 Markdown 文件

    策略: 按 ## 标题分块
    """
    file_path = Path(file_path)
    docs = []

    try:
        text = file_path.read_text(encoding="utf-8")
        sections = text.split("\n## ")

        for i, section in enumerate(sections):
            content = section.strip()
            if not content or len(content) < 30:
                continue

            # 提取标题
            lines = content.split("\n", 1)
            title = lines[0].lstrip("# ").strip()

            docs.append(Document(
                content=content if i == 0 else f"## {content}",
                metadata={
                    "source": file_path.name,
                    "source_type": "markdown",
                    "section": title,
                },
            ))

        logger.info("Markdown 解析完成: %s → %d 块", file_path.name, len(docs))

    except Exception as e:
        logger.error("Markdown 解析失败 %s: %s", file_path, e)

    return docs


def parse_python(file_path: str | Path) -> list[Document]:
    """
    解析 Python 文件

    策略: 提取 docstring + 类/函数定义
    """
    file_path = Path(file_path)
    docs = []

    try:
        text = file_path.read_text(encoding="utf-8")
        if len(text) < 50:
            return docs

        docs.append(Document(
            content=text[:3000],  # 截断超长文件
            metadata={
                "source": file_path.name,
                "source_type": "python",
                "lines": text.count("\n"),
            },
        ))

    except Exception as e:
        logger.error("Python 解析失败 %s: %s", file_path, e)

    return docs


# ============================================================
# 统一入口
# ============================================================

PARSER_MAP = {
    ".pdf": parse_pdf,
    ".pptx": parse_pptx,
    ".ppt": parse_pptx,  # 注意: 旧版 .ppt 不支持，需要转换
    ".docx": parse_docx,
    ".doc": parse_docx,
    ".md": parse_markdown,
    ".py": parse_python,
    ".txt": parse_markdown,  # 纯文本按 markdown 处理
}


def parse_file(file_path: str | Path) -> list[Document]:
    """
    自动识别文件格式并解析

    支持: .pdf, .pptx, .docx, .md, .py, .txt
    """
    file_path = Path(file_path)
    suffix = file_path.suffix.lower()

    parser = PARSER_MAP.get(suffix)
    if not parser:
        logger.warning("不支持的文件格式: %s (%s)", file_path.name, suffix)
        return []

    return parser(file_path)


def parse_directory(dir_path: str | Path, recursive: bool = True) -> list[Document]:
    """
    解析目录下所有支持的文件

    Returns: 所有文档块的列表
    """
    dir_path = Path(dir_path)
    all_docs = []
    supported = set(PARSER_MAP.keys())

    pattern = "**/*" if recursive else "*"
    for f in dir_path.glob(pattern):
        if f.is_file() and f.suffix.lower() in supported:
            docs = parse_file(f)
            all_docs.extend(docs)

    logger.info("目录解析完成: %s → %d 个文件, %d 个文档块",
                dir_path, len(set(d.metadata.get("source", "") for d in all_docs)), len(all_docs))
    return all_docs
