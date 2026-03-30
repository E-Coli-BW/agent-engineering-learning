"""
文档解析器测试
===============
测试 PDF/PPT/DOCX/Markdown/Python 解析器
使用项目自身的文件做测试数据
"""

import tempfile
from pathlib import Path


class TestMarkdownParser:
    def test_parse_readme(self):
        from project.etl.parsers import parse_markdown
        readme = Path(__file__).parent.parent / "README.md"
        docs = parse_markdown(readme)
        assert len(docs) > 0
        assert docs[0].metadata["source_type"] == "markdown"
        assert docs[0].metadata["source"] == "README.md"

    def test_parse_with_sections(self):
        from project.etl.parsers import parse_markdown
        # 创建临时 md 文件
        with tempfile.NamedTemporaryFile(suffix=".md", mode="w", delete=False) as f:
            f.write("# Title\n\nIntro paragraph here.\n\n## Section A\n\nContent of section A.\n\n## Section B\n\nContent of section B.\n")
            f.flush()
            docs = parse_markdown(f.name)
            # 应该有多个 section
            assert len(docs) >= 2


class TestPythonParser:
    def test_parse_python_file(self):
        from project.etl.parsers import parse_python
        config = Path(__file__).parent.parent / "project" / "config.py"
        docs = parse_python(config)
        assert len(docs) == 1
        assert docs[0].metadata["source_type"] == "python"
        assert "Config" in docs[0].content


class TestDocxParser:
    def test_parse_generated_docx(self):
        """生成一个临时 DOCX 并解析"""
        from docx import Document as DocxDoc
        from project.etl.parsers import parse_docx

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc = DocxDoc()
            doc.add_heading("培训手册", level=1)
            doc.add_paragraph("本手册介绍呼叫中心坐席的标准操作流程。")
            doc.add_heading("投诉处理 SOP", level=2)
            doc.add_paragraph("第一步：倾听客户诉求。\n第二步：确认问题。\n第三步：提供解决方案。")
            doc.add_heading("常见问题", level=2)
            doc.add_paragraph("Q: 客户要求转接主管怎么办？\nA: 按照转接流程操作。")
            doc.save(f.name)

            docs = parse_docx(f.name)
            assert len(docs) >= 2
            assert any("投诉处理" in d.metadata.get("section", "") for d in docs)
            assert docs[0].metadata["source_type"] == "docx"


class TestPptxParser:
    def test_parse_generated_pptx(self):
        """生成一个临时 PPTX 并解析"""
        from pptx import Presentation
        from pptx.util import Inches
        from project.etl.parsers import parse_pptx

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()

            # Slide 1: 标题页
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "呼叫中心培训"
            slide.placeholders[1].text = "2026年新员工入职培训"

            # Slide 2: 内容页
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "投诉处理流程"
            body = slide.placeholders[1]
            body.text = "1. 倾听客户\n2. 确认问题\n3. 提供方案\n4. 跟进反馈"

            # Slide 3: 带备注
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "常见话术"
            slide.placeholders[1].text = "开场白：您好，很高兴为您服务。"
            notes = slide.notes_slide
            notes.notes_text_frame.text = "注意语气要温和，不要带情绪"

            prs.save(f.name)

            docs = parse_pptx(f.name)
            assert len(docs) == 3  # 3 slides
            assert docs[0].metadata["source_type"] == "pptx"
            assert docs[0].metadata["slide"] == 1
            assert "呼叫中心培训" in docs[0].content
            # 第三张 slide 应该有备注
            assert docs[2].metadata["has_notes"] is True
            assert "语气" in docs[2].content


class TestParseFile:
    def test_auto_detect_markdown(self):
        from project.etl.parsers import parse_file
        readme = Path(__file__).parent.parent / "README.md"
        docs = parse_file(readme)
        assert len(docs) > 0

    def test_unsupported_format(self):
        from project.etl.parsers import parse_file
        with tempfile.NamedTemporaryFile(suffix=".xyz", delete=False) as f:
            f.write(b"test")
            f.flush()
            docs = parse_file(f.name)
            assert docs == []


class TestParseDirectory:
    def test_parse_project_dir(self):
        from project.etl.parsers import parse_directory
        # 解析 agent/ 目录 (只有 .py 和 .md)
        agent_dir = Path(__file__).parent.parent / "agent"
        docs = parse_directory(agent_dir, recursive=False)
        assert len(docs) > 0
        source_types = set(d.metadata["source_type"] for d in docs)
        assert "python" in source_types or "markdown" in source_types
